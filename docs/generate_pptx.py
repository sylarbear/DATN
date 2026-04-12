# -*- coding: utf-8 -*-
"""
Generate graduation thesis presentation (PowerPoint)
English Learning - Phan Quang Thuật
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DOCS_DIR = r"g:\xamp\htdocs\DATN\docs"
SCREENSHOTS_DIR = os.path.join(DOCS_DIR, "screenshots")
UML_DIR = os.path.join(DOCS_DIR, "uml_images")
OUTPUT = os.path.join(DOCS_DIR, "slide_bao_ve.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============================================================
# COLOR PALETTE
# ============================================================
BLUE_DARK = RGBColor(0x1e, 0x3a, 0x5f)
BLUE_PRIMARY = RGBColor(0x3b, 0x5b, 0x9d)
BLUE_LIGHT = RGBColor(0x6c, 0x8e, 0xbf)
BLUE_ACCENT = RGBColor(0x45, 0x6b, 0xe8)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xf0, 0xf2, 0xf5)
DARK_TEXT = RGBColor(0x2d, 0x2d, 0x2d)
MID_TEXT = RGBColor(0x55, 0x55, 0x55)
ORANGE = RGBColor(0xf5, 0x9e, 0x0b)
GREEN = RGBColor(0x10, 0xb9, 0x81)
RED_SOFT = RGBColor(0xef, 0x44, 0x44)

# ============================================================
# HELPER FUNCTIONS
# ============================================================
def add_bg_rect(slide, color=BLUE_DARK):
    """Add full-width background rectangle"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_bg_gradient(slide, color1=BLUE_DARK, color2=BLUE_PRIMARY):
    """Add gradient background"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = color2
    fill.gradient_stops[1].position = 1.0
    shape.line.fill.background()

def add_header_bar(slide, title_text):
    """Add blue header bar at top"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE_DARK
    bar.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(0.85))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"

def add_accent_line(slide, left, top, width):
    """Add decorative accent line"""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.fill.background()

def add_text_box(slide, left, top, width, height, text, size=18, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT, font_name="Arial"):
    """Add a text box"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf

def add_bullet_list(slide, left, top, width, height, items, size=16, color=DARK_TEXT, spacing=1.2):
    """Add bullet list"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.space_after = Pt(size * spacing * 0.5)
    return tf

def add_image_safe(slide, img_path, left, top, width=None, height=None):
    """Add image safely"""
    if os.path.exists(img_path):
        kwargs = {}
        if width: kwargs['width'] = width
        if height: kwargs['height'] = height
        slide.shapes.add_picture(img_path, left, top, **kwargs)
        return True
    return False

def add_number_box(slide, left, top, number, label, box_color=BLUE_ACCENT):
    """Add a stat box with number and label"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.5), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = box_color
    box.line.fill.background()
    
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(number)
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Arial"
    
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(13)
    run2.font.color.rgb = RGBColor(0xe0, 0xe8, 0xf0)
    run2.font.name = "Arial"

# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg_gradient(slide, BLUE_DARK, RGBColor(0x0f, 0x25, 0x45))

# Accent line
add_accent_line(slide, Inches(4), Inches(1.8), Inches(5.3))

# School name
add_text_box(slide, Inches(1), Inches(0.5), Inches(11.3), Inches(0.6),
    "TRƯỜNG CAO ĐẲNG CÔNG THƯƠNG TPHCM – KHOA CNTT", 16, False, BLUE_LIGHT, PP_ALIGN.CENTER)

# Title
add_text_box(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.2),
    "XÂY DỰNG WEBSITE HỌC TIẾNG ANH\nTRỰC TUYẾN ENGLISH LEARNING", 40, True, WHITE, PP_ALIGN.CENTER)

# Accent line bottom
add_accent_line(slide, Inches(4), Inches(3.6), Inches(5.3))

# Info
add_text_box(slide, Inches(2.5), Inches(4.2), Inches(4), Inches(1.5),
    "GVHD: Vũ Thị Hường", 20, False, BLUE_LIGHT, PP_ALIGN.LEFT)
add_text_box(slide, Inches(7), Inches(4.2), Inches(4), Inches(1.5),
    "SVTH: Phan Quang Thuật\nMSSV: 2120110351\nLớp: CCQ2011E – Khóa K44", 20, False, BLUE_LIGHT, PP_ALIGN.LEFT)

add_text_box(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.5),
    "BÁO CÁO THỰC TẬP TỐT NGHIỆP • TPHCM, 04/2026", 14, False, RGBColor(0x80,0x90,0xa8), PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: NỘI DUNG TRÌNH BÀY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "📋 NỘI DUNG TRÌNH BÀY")

items = [
    ("01", "Giới thiệu đề tài", "Lý do chọn đề tài, mục tiêu, phạm vi"),
    ("02", "Cơ sở lý thuyết", "Công nghệ sử dụng: PHP, MySQL, MVC, API"),
    ("03", "Phân tích thiết kế", "UML: Use Case, Sequence, State Diagram"),
    ("04", "Demo sản phẩm", "Giao diện người dùng và quản trị"),
    ("05", "Kết luận", "Kết quả, hạn chế, hướng phát triển"),
]

for i, (num, title, desc) in enumerate(items):
    y = Inches(1.5) + Inches(i * 1.1)
    
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y, Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BLUE_ACCENT
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = num
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = WHITE

    add_text_box(slide, Inches(2.2), y + Inches(0.02), Inches(4), Inches(0.4), title, 22, True, DARK_TEXT)
    add_text_box(slide, Inches(2.2), y + Inches(0.42), Inches(8), Inches(0.35), desc, 15, False, MID_TEXT)

# ============================================================
# SLIDE 3: GIỚI THIỆU ĐỀ TÀI
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "01 — GIỚI THIỆU ĐỀ TÀI")

add_text_box(slide, Inches(0.6), Inches(1.4), Inches(5.8), Inches(0.5),
    "Lý do chọn đề tài", 24, True, BLUE_DARK)

add_bullet_list(slide, Inches(0.6), Inches(2.0), Inches(5.8), Inches(4.5), [
    "Nhu cầu học tiếng Anh online tại Việt Nam tăng mạnh",
    "AI (GPT, Speech API) cho phép đánh giá Speaking tự động",
    "Tích hợp đa kỹ năng: Nghe, Nói, Đọc, từ vựng, ngữ pháp",
    "Hệ thống ví điện tử + membership chuyên nghiệp",
    "Áp dụng toàn diện kiến thức: PHP MVC, MySQL, JS, API",
], 16, DARK_TEXT)

# Right side - Homepage screenshot
add_image_safe(slide, os.path.join(SCREENSHOTS_DIR, "02_homepage_top.png"),
    Inches(7), Inches(1.4), width=Inches(5.8))

# ============================================================
# SLIDE 4: MỤC TIÊU ĐỀ TÀI
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "01 — MỤC TIÊU ĐỀ TÀI")

add_text_box(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(0.5),
    "Mục tiêu cụ thể", 24, True, BLUE_DARK)

# Left column
add_bullet_list(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(4.5), [
    "Xây dựng hệ thống học từ vựng, ngữ pháp, bài kiểm tra",
    "Tích hợp AI đánh giá kỹ năng nói (Speaking)",
    "Phát triển tính năng Flashcard học từ vựng",
    "Xây dựng bài kiểm tra: Quiz, Listening, Reading",
], 16)

# Right column
add_bullet_list(slide, Inches(6.8), Inches(2.1), Inches(5.8), Inches(4.5), [
    "Phát triển hệ thống ví điện tử + membership",
    "Dashboard thống kê tiến độ + Leaderboard",
    "Hệ thống XP, Level, Streak (Gamification)",
    "Trang quản trị Admin đầy đủ CRUD + tài chính",
], 16)

# Stats boxes
add_number_box(slide, Inches(0.8), Inches(5.5), "23", "Bảng CSDL", BLUE_ACCENT)
add_number_box(slide, Inches(3.8), Inches(5.5), "30+", "Trang giao diện", RGBColor(0x10,0xb9,0x81))
add_number_box(slide, Inches(6.8), Inches(5.5), "6", "API tích hợp", ORANGE)
add_number_box(slide, Inches(9.8), Inches(5.5), "4", "Kỹ năng tiếng Anh", RGBColor(0xef,0x44,0x44))

# ============================================================
# SLIDE 5: CÔNG NGHỆ SỬ DỤNG
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "02 — CÔNG NGHỆ SỬ DỤNG")

techs = [
    ("Backend", ["PHP 8.2 (MVC thuần)", "MySQL + PDO", "Session + Cookie"], BLUE_ACCENT),
    ("Frontend", ["HTML5 + CSS3", "JavaScript (AJAX)", "Chart.js (biểu đồ)"], GREEN),
    ("API & Tích hợp", ["Google OAuth 2.0", "OpenAI GPT API", "Web Speech API"], ORANGE),
    ("Thanh toán", ["VietQR (QR Code)", "Ví điện tử", "Mã kích hoạt"], RGBColor(0xef, 0x44, 0x44)),
]

for i, (title, items, color) in enumerate(techs):
    x = Inches(0.5) + Inches(i * 3.15)
    y = Inches(1.6)
    
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.9), Inches(4.8))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xe0, 0xe0, 0xe0)
    card.line.width = Pt(1)
    
    # Color bar on top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(2.9), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    
    # Title
    add_text_box(slide, x + Inches(0.2), y + Inches(0.3), Inches(2.5), Inches(0.5),
        title, 20, True, color, PP_ALIGN.CENTER)
    
    # Items
    for j, item in enumerate(items):
        item_y = y + Inches(1.1) + Inches(j * 1.1)
        # Bullet dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.4), item_y + Inches(0.08), Inches(0.15), Inches(0.15))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        
        add_text_box(slide, x + Inches(0.7), item_y, Inches(2.0), Inches(0.9),
            item, 15, False, DARK_TEXT, PP_ALIGN.LEFT)

# ============================================================
# SLIDE 6: MÔ HÌNH MVC
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "02 — MÔ HÌNH MVC")

# MVC boxes
mvc = [
    ("Model", "Xử lý logic nghiệp vụ\nTương tác CSDL\n\nVí dụ: User, Topic,\nTest, Vocabulary...", BLUE_ACCENT, Inches(0.8)),
    ("View", "Giao diện người dùng\nHiển thị dữ liệu\n\nVí dụ: home/index,\nlessons/show...", GREEN, Inches(4.7)),
    ("Controller", "Tiếp nhận request\nĐiều hướng logic\n\nVí dụ: AuthController,\nWalletController...", ORANGE, Inches(8.6)),
]

for title, desc, color, x in mvc:
    # Box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.5), Inches(4.2))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(2)
    
    # Header
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.8), Inches(3.5), Inches(0.8))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = color
    hdr.line.fill.background()
    
    add_text_box(slide, x + Inches(0.2), Inches(1.9), Inches(3.1), Inches(0.6),
        title, 26, True, WHITE, PP_ALIGN.CENTER)
    
    add_text_box(slide, x + Inches(0.3), Inches(2.9), Inches(2.9), Inches(2.8),
        desc, 16, False, DARK_TEXT, PP_ALIGN.CENTER)

# Structure info
add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.8),
    "📁 Cấu trúc: app/controllers/ → app/models/ → app/views/ → public/", 15, False, MID_TEXT, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 7: USE CASE DIAGRAM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "03 — SƠ ĐỒ USE CASE")

add_image_safe(slide, os.path.join(UML_DIR, "use_case.png"),
    Inches(0.5), Inches(1.3), width=Inches(7))

# Actor list
add_text_box(slide, Inches(8), Inches(1.5), Inches(4.5), Inches(0.4), "Actors:", 20, True, BLUE_DARK)

actors = [
    ("👤 Guest", "Xem chủ đề, đăng ký, đăng nhập"),
    ("👨‍🎓 User (Free)", "Học từ, quiz, flashcard, ví, bookmark"),
    ("⭐ Pro User", "+ Listening, Reading, Speaking"),
    ("🔧 Admin", "+ Quản lý hệ thống, duyệt đơn"),
]
for i, (actor, desc) in enumerate(actors):
    y = Inches(2.2) + Inches(i * 1.1)
    add_text_box(slide, Inches(8), y, Inches(4.5), Inches(0.35), actor, 17, True, BLUE_ACCENT)
    add_text_box(slide, Inches(8), y + Inches(0.35), Inches(4.5), Inches(0.5), desc, 14, False, MID_TEXT)

# ============================================================
# SLIDE 8: SEQUENCE DIAGRAMS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "03 — SƠ ĐỒ SEQUENCE: LUỒNG CHÍNH")

# 2x2 grid
seqs = [
    ("Đăng nhập", "sequence_login.png"),
    ("Nạp tiền vào ví", "sequence_deposit.png"),
    ("Mua gói Pro", "sequence_purchase.png"),
    ("Hủy đơn + Hoàn tiền", "sequence_cancel.png"),
]

positions = [
    (Inches(0.3), Inches(1.3)), (Inches(6.8), Inches(1.3)),
    (Inches(0.3), Inches(4.3)), (Inches(6.8), Inches(4.3)),
]

for (title, img), (x, y) in zip(seqs, positions):
    add_text_box(slide, x, y, Inches(6), Inches(0.35), f"▸ {title}", 14, True, BLUE_DARK)
    add_image_safe(slide, os.path.join(UML_DIR, img), x, y + Inches(0.35), width=Inches(6))

# ============================================================
# SLIDE 9: CƠ SỞ DỮ LIỆU
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "03 — THIẾT KẾ CƠ SỞ DỮ LIỆU")

add_text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
    "MySQL — 23 bảng — Charset utf8mb4", 20, True, BLUE_DARK)

# Table groups
groups = [
    ("Học tập", ["topics", "vocabularies", "lessons", "lesson_contents", "tests", "questions", "grammar_lessons", "grammar_questions", "speaking_prompts", "speaking_attempts"], BLUE_ACCENT),
    ("Người dùng", ["users", "user_progress", "bookmarks", "lesson_reviews", "test_results", "user_answers", "xp_history"], GREEN),
    ("Tài chính", ["membership_plans", "membership_orders", "wallet_transactions", "activation_codes"], ORANGE),
    ("Hỗ trợ", ["support_tickets", "vocab_reviews"], RED_SOFT),
]

for i, (group_name, tables, color) in enumerate(groups):
    x = Inches(0.5) + Inches(i * 3.15)
    
    # Group header
    hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0), Inches(2.9), Inches(0.55))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = color
    hdr.line.fill.background()
    add_text_box(slide, x + Inches(0.1), Inches(2.05), Inches(2.7), Inches(0.45),
        f"{group_name} ({len(tables)})", 16, True, WHITE, PP_ALIGN.CENTER)
    
    # Table list
    table_text = "\n".join([f"• {t}" for t in tables])
    add_text_box(slide, x + Inches(0.2), Inches(2.7), Inches(2.7), Inches(4.2),
        table_text, 12, False, DARK_TEXT)

# ============================================================
# SLIDE 10-16: DEMO SCREENSHOTS
# ============================================================
demo_pages = [
    ("04 — TRANG CHỦ & ĐĂNG NHẬP",
     [("02_homepage_top.png", "Trang chủ"), ("01_login.png", "Đăng nhập")]),
    ("04 — HỌC TỪ VỰNG & FLASHCARD",
     [("05_topic_detail.png", "Chi tiết chủ đề + từ vựng"), ("07_flashcard.png", "Flashcard")]),
    ("04 — BÀI KIỂM TRA & NGỮ PHÁP",
     [("09_test.png", "Bài kiểm tra"), ("08_grammar.png", "Ngữ pháp")]),
    ("04 — LUYỆN NÓI & DASHBOARD",
     [("10_speaking.png", "Luyện nói (AI)"), ("11_dashboard.png", "Dashboard tiến độ")]),
    ("04 — VÍ ĐIỆN TỬ & NẠP TIỀN",
     [("15_wallet.png", "Ví điện tử"), ("16_deposit.png", "Nạp tiền QR")]),
    ("04 — NÂNG CẤP PRO & HỖ TRỢ",
     [("18_membership_pricing.png", "Bảng giá Pro"), ("20_support_list.png", "Tickets hỗ trợ")]),
    ("04 — QUẢN TRỊ ADMIN",
     [("22_admin_dashboard.png", "Admin Dashboard"), ("29_admin_wallet.png", "Quản lý ví")]),
]

for title, images in demo_pages:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, title)
    
    for j, (img_file, caption) in enumerate(images):
        x = Inches(0.3) + Inches(j * 6.5)
        add_text_box(slide, x + Inches(0.2), Inches(1.25), Inches(6), Inches(0.4),
            caption, 16, True, BLUE_DARK, PP_ALIGN.CENTER)
        add_image_safe(slide, os.path.join(SCREENSHOTS_DIR, img_file),
            x + Inches(0.1), Inches(1.7), width=Inches(6))

# ============================================================
# SLIDE 17: STATE DIAGRAMS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "03 — SƠ ĐỒ TRẠNG THÁI")

states = [
    ("Đơn hàng", "state_order.png"),
    ("Giao dịch ví", "state_wallet.png"),
    ("Ticket hỗ trợ", "state_ticket.png"),
]

for i, (title, img) in enumerate(states):
    x = Inches(0.3) + Inches(i * 4.3)
    add_text_box(slide, x, Inches(1.4), Inches(4), Inches(0.4), f"▸ {title}", 16, True, BLUE_DARK, PP_ALIGN.CENTER)
    add_image_safe(slide, os.path.join(UML_DIR, img), x + Inches(0.2), Inches(1.9), width=Inches(3.7))

# ============================================================
# SLIDE 18: KẾT QUẢ KIỂM THỬ
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "05 — KẾT QUẢ KIỂM THỬ")

add_number_box(slide, Inches(0.8), Inches(1.8), "30+", "Trang đã test", BLUE_ACCENT)
add_number_box(slide, Inches(3.8), Inches(1.8), "0", "Lỗi PHP", GREEN)
add_number_box(slide, Inches(6.8), Inches(1.8), "100%", "Encoding OK", ORANGE)
add_number_box(slide, Inches(9.8), Inches(1.8), "✓", "Wallet Flow", RGBColor(0x10,0xb9,0x81))

add_text_box(slide, Inches(0.6), Inches(3.5), Inches(5.8), Inches(0.4),
    "Các luồng đã kiểm thử:", 20, True, BLUE_DARK)

add_bullet_list(slide, Inches(0.6), Inches(4.1), Inches(5.5), Inches(3), [
    "Đăng ký, đăng nhập, Google OAuth",
    "Học từ vựng, flashcard, bookmark",
    "Làm bài quiz, listening, reading",
    "Luyện nói + AI chấm điểm",
    "Nạp tiền → Admin duyệt → Mua gói Pro",
    "Hủy đơn → Admin duyệt → Hoàn tiền vào ví",
    "Rút tiền → Admin chuyển khoản",
], 15)

add_bullet_list(slide, Inches(7), Inches(4.1), Inches(5.5), Inches(3), [
    "Dashboard thống kê tiến độ",
    "Bảng xếp hạng XP",
    "Admin: CRUD chủ đề, câu hỏi, user",
    "Admin: Duyệt đơn, quản lý ví",
    "Admin: Phản hồi tickets",
    "Vietnamese encoding (UTF-8)",
    "Responsive design (PC + Mobile)",
], 15)

# ============================================================
# SLIDE 19: KẾT LUẬN
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "05 — KẾT LUẬN")

# Ưu điểm
add_text_box(slide, Inches(0.6), Inches(1.4), Inches(5.8), Inches(0.4),
    "✅ Ưu điểm", 22, True, GREEN)
add_bullet_list(slide, Inches(0.6), Inches(1.9), Inches(5.8), Inches(2.5), [
    "Hoàn thành đầy đủ chức năng đề ra",
    "Tích hợp AI đánh giá Speaking – điểm nổi bật",
    "Hệ thống ví + membership chuyên nghiệp",
    "Giao diện đẹp, hiện đại, responsive",
    "Gamification: XP, Level, Streak, Leaderboard",
    "Bảo mật tài chính: transaction + row locking",
], 14, DARK_TEXT)

# Hạn chế
add_text_box(slide, Inches(7), Inches(1.4), Inches(5.8), Inches(0.4),
    "⚠️ Hạn chế", 22, True, ORANGE)
add_bullet_list(slide, Inches(7), Inches(1.9), Inches(5.8), Inches(2.5), [
    "Nạp tiền chưa tự động (cần Admin duyệt)",
    "Chưa có app mobile native",
    "Nội dung còn hạn chế (6 chủ đề)",
    "Speaking chỉ trên Chrome/Edge",
], 14, DARK_TEXT)

# Hướng phát triển
add_text_box(slide, Inches(0.6), Inches(4.8), Inches(12), Inches(0.4),
    "🚀 Hướng phát triển", 22, True, BLUE_ACCENT)
add_bullet_list(slide, Inches(0.6), Inches(5.4), Inches(5.8), Inches(2), [
    "Tích hợp VNPay, Momo thanh toán tự động",
    "Phát triển app mobile (React Native/Flutter)",
    "Bổ sung nhiều chủ đề, video bài giảng",
], 14, DARK_TEXT)
add_bullet_list(slide, Inches(7), Inches(5.4), Inches(5.8), Inches(2), [
    "Tích hợp Spaced Repetition (SRS)",
    "Thêm forum, chat, học nhóm",
    "Deploy lên hosting thật (Cloud/VPS)",
], 14, DARK_TEXT)

# ============================================================
# SLIDE 20: CẢM ƠN
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_gradient(slide, BLUE_DARK, RGBColor(0x0f, 0x25, 0x45))

add_accent_line(slide, Inches(4.5), Inches(2.5), Inches(4.3))

add_text_box(slide, Inches(1), Inches(2.8), Inches(11.3), Inches(1.2),
    "CẢM ƠN QUÝ THẦY CÔ\nĐÃ LẮNG NGHE!", 42, True, WHITE, PP_ALIGN.CENTER)

add_accent_line(slide, Inches(4.5), Inches(4.5), Inches(4.3))

add_text_box(slide, Inches(1), Inches(5.2), Inches(11.3), Inches(0.5),
    "SVTH: Phan Quang Thuật • MSSV: 2120110351 • GVHD: Vũ Thị Hường", 18, False, BLUE_LIGHT, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.8), Inches(11.3), Inches(0.5),
    "GitHub: github.com/sylarbear/DATN", 16, False, RGBColor(0x80,0x90,0xa8), PP_ALIGN.CENTER)

# ============================================================
# SAVE
# ============================================================
prs.save(OUTPUT)
print(f"Done! Saved to: {OUTPUT}")
print(f"File size: {os.path.getsize(OUTPUT) / 1024:.0f} KB")
print(f"Total slides: {len(prs.slides)}")
